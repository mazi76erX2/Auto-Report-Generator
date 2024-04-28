import random
import time
from datetime import datetime
from urllib.error import HTTPError, URLErrors
from urllib.parse import quote_plus
from urllib.request import urlopen

import pandas as pd
import requests
from bs4 import BeautifulSoup
from openpyxl import load_workbook
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

LIST_PRINT_TOP10 = []
LIST_ONLINE_TOP10 = []

BITLY_API = 'e8c05d9133d039c84ab05395e7871e0c0e3b2432'

SUBJECT_LIST = ['GT-R', 'NV 350 Impendulo', 'Qashqai', 'NV 200 Combi',
                'Tiida', 'Leaf', 'NAAMSA', 'K-line', 'Juke', 'X-Trail',
                'Navara', 'Micra', 'NP 200', 'Patrol', 'Sentra',
                'Nissan Festival of Motoring',
                'Nissan drives for further growth',
                'ICC Sponsorship', 'NP 300', 'Nissan Trailseeker',
                'Nissan Corporate',
                'Nissan Taxi Call Centre oppurtunity',
                'Nissn long distance driving tips',
                'Nissan Holiday Checklist', 'Pathfinder', 'NP 300 hardbody',
                'NV 300', 'Workhorse', 'Almera', '350z', 'Hardbody',
                'Simola Hillclimb', '370Z', 'Murano',
                'Nissan Kicks', 'Nissan Primera', 'Bladeglider',
                'Nissan ProPilot', 'Festival of Motoring', 'NV 350 Panel Van',
                'Patrol Wagon', 'Livina', 'Nissan Sponsorship/sport',
                'Nissan Motorsport', 'Nissan Spokespeople',
                'General Industry related', 'Titan', 'Rouge']


def loadPrime():
    primeList = ('data-files/Prime Reading Current as of 02.2018.xlsx')

    df = pd.read_excel(primeList)

    printList = list(df['PRINT '])
    printList = list(filter(lambda x: x == x, printList))

    onlineList = list(df['ONLINE '])
    onlineList = list(filter(lambda x: x == x, onlineList))

    tvList = list(df['TV'])
    tvList = list(filter(lambda x: x == x, tvList))

    return printList, onlineList, tvList


def findSubject(text, description):
    productFound = ''
    for product in SUBJECT_LIST:
        if (text.lower().find(product.lower()) != -1) or \
           (description.lower().find(product.lower()) != -1):
            productFound = product

    return productFound


def findPrimePrint(media):
    primeString = ''
    for prime in PRINT_LIST:
        if media.lower().find(prime.lower()) != -1:
            primeString = ' [PRIME]'

    return primeString


def findPrimeOnline(media):
    primeString = ''
    for prime in ONLINE_LIST:
        if media.lower().find(prime.lower()) != -1:
            primeString = ' [PRIME]'

    return primeString


def findPrimeBroadcast(media):
    primeString = ''
    for prime in TV_LIST:
        if media.lower().find(prime.lower()) != -1:
            primeString = ' [PRIME]'

    return primeString


def followLink(tweet):
    """Obtains an extract from a tweet if one is available
    """
    pass


def checkIfPrime(media):
    """Check if it's prime.
    """
    if media.find == 'PRIME':
        return True
    else:
        return False


def getExtract(url, description):
    text = str(description)

    try:
        page = urlopen(url)
        soup = BeautifulSoup(page, 'html.parser')
        paragraphs = soup.find_all('p')

        for paragraph in paragraphs:
            text += paragraph.get_text() + ' '

        while len(text) > 700:
            index = text.rfind('.')
            text = text[:index]

        text += '.'

     except (HTTPError, URLError) as error:
        text = description
        
    return text


def set_table_styles(table):
    """styling from Series to the table as a list
    """
    for cell in table:
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        cell.text_frame.paragraphs[0].font.name = 'Arial'
        cell.text_frame.paragraphs[0].font.bold = True

    return table
    

def set_table_styles_link_and_extract(table):
    ##  Not working
    ##  http://python-pptx.readthedocs.io/en/latest/dev/analysis/txt-hyperlink.html
    #  Styling for link
    table.cell(8,2).text_frame.paragraphs[0].add_run()
    table.cell(8,2).text_frame.paragraphs[0].hyperlink =\
                                            'https://www.google.com'
    table.cell(8,2).text_frame.paragraphs[0].font.underline = True

    #  Styling for extract
    table.cell(9,2).text_frame.paragraphs[0].font.bold = False
    table.cell(9,2).text_frame.paragraphs[0].font.italic = True
    table.cell(9,2).text_frame.paragraphs[0].font.color.rgb =\
                                                            RGBColor(89,89,89)

    return table


def top10Stories(wb, ws):
    """Generates Top10 from print and online clips
    """
    #FIX!!!!!!!!!!!!!!!!!!
    ws = ws['J11':'L20']

    #Randomly selects 5 print clips for print and online
    selection1 = random.sample(range(0, 10), 5)
    selection2 = random.sample(range(0, 10), 5)

    for index in range(5):
        if checkIfPrime(LIST_PRINT_TOP10[selection1[index]][0]):
            ws[index][2].value = 'PRIME'

        media = LIST_PRINT_TOP10[selection1[index]][0].replace(' [PRIME]', '')
        
        ws[index][0].value = media                                    #media
        ws[index][1].value = LIST_PRINT_TOP10[selection1[index]][1]   #title

    for index in range(5,10):
        if checkIfPrime(LIST_ONLINE_TOP10[selection2[index-5]][0]):
            ws[index][2].value = 'PRIME'

        media = LIST_ONLINE_TOP10[selection2[index-5]][0].replace(' [PRIME]',
                                                                  '')
        
        ws[index][0].value = media                                     #media
        ws[index][1].value = LIST_ONLINE_TOP10[selection2[index-5]][1] #title

    wb.save('done.xlsx')
 

def populate_table_to_list(table):
    """Populate list from the Table (pptx.shapes.table object) values into
        a list of Cells (pptx.shapes.table.cell objects)
    """
    cell_list = []

    for i in range(1,10):
        cell_list.append(table.cell(i,2))

    return cell_list


def printRow(row, coverage):
    #Converts the url
    url = quote_plus(row[3].value)
    print('URL working: ', row[3].value)
    
    #Bitly GET call address
    api_call = 'https://api-ssl.bitly.com/v3/shorten?access_token={}'\
           '&longUrl={}'.format(BITLY_API, url)
    
    
    #Using requests for the GET call
    url = requests.get(api_call).json()['data']['url']
    print('API working')

    #These variables are added for readability

    title = row[0].value
    media = row[1].value #+ findPrimePrint(row[1].value)
    date = row[11].value
    tone = row[8].value
    reach = row[9].value
    ave = row[7].value
    link = url
    extract = row[4].value.replace('_x000D_\n',' ').replace('\ufeff','')

    count = 0
    while len(extract)-1 > 500:
        index = extract.rfind('.')
        extract = extract[:index]
        count += 1

    extract = extract + '.'

    if coverage == 'Product':
        nameplate = row[12].value
    else:
        nameplate = findSubject(extract, row[5].value)

    nameplate = nameplate.replace(', Nissan', '')

    #Might use the code below in the future
    #commaIndex = row[16].value.find(',')
    #subjects = row[16].value[:4]

    return [nameplate, title, media, date, tone, reach, ave, link, extract]


def onlineRow(row, coverage):
    """
    parses row of sheet
    """
    #Converts the url
    url = quote_plus(row[4].value)

    #Bitly GET call address
    api_call = 'https://api-ssl.bitly.com/v3/shorten?access_token={}'\
           '&longUrl={}'.format(BITLY_API, url)

    #Using requests for the GET call
    url = requests.get(api_call).json()['data']['url'] 

    #These variables are added for readability
    title = row[2].value
    

    if coverage == 'SSA Online':
        media = row[12].value
        index = media.find('//')
        media = media[index+2:]
        media = media.replace('www.','')
        media = media #+ findPrimeOnline(media) 
        tone = row[10].value
        reach = row[18].value
        ave = row[19].value
        #Get Longer extract of description
        extract = getExtract(row[4].value, row[3].value)
        if extract == '':
            row[3].value
    elif coverage == 'Social':
        media = row[11].value
        tone = row[11].value
        reach = row[17].value
        ave = row[19].value
        extract = row[2].value
        index = title.find('http')
        title = extract[:index-1]
    else: #  Online
        media = row[13].value
        index = media.find('//')
        media = media[index+2:]
        media = media.replace('www.','')
        media = media #+ findPrimeOnline(media) 
        tone = row[11].value
        reach = row[19].value
        ave = row[20].value
        extract = getExtract(row[4].value, row[3].value)
        if len(extract) < 100:
            row[3].value
    
    date = row[5].value[:10]

    nameplate = findSubject(extract, row[2].value) 

    link = url

    return [nameplate, title, media, date, tone.capitalize(), reach, ave,
            link, extract]


def determine_slide(coverage, media):
    """
    """
    # Get reference to second slide
    slide = prs.slides[3]
    table = sl.shapes[3].table

    return table


def (prs):
    print_product_slides = [3,4,5,6,7]
    print_coverage_slides = [15,16,17,18,19]

    online_product_slides = [9,10,11,12,13]
    online_coverage_slides = [21,22,23,24,25]

    broadcast_slides = [27,28]

    online_ssa_slides = [30,31,32,33,34]
    social_ssa_slides = [36,37,38]


    

    table = set_table_values(wb, wsPrint, table, 'Product', 'Print')


def set_table_values_old(wb, ws1, table, coverage, media):
    print('')
    #Colours of the highlighted cells
    colours = {'Product' : 'FFFFFF00',  'Corporate' : 'FF92D050',
               'SSA Online': 'FFFFC000', 'Social': 'FF00B0F0'}
    colour = colours[coverage]
    
    #creates list of rows from ws1
    rows = list(ws1.rows)

    #Increases by 11 each loop because of the excel sheets setup 
    clip_count = 0

    #Count the number of clips added
    count = 0 

    #iterate first cell in each row until you find a yellow cell 
    for row in rows:
        #If the cell is a certan colour
        if row[0].fill.fgColor.rgb == colour:
            print(row[0].value)

            #Max 5 clips
            if (count==3) and (coverage=='Social'):
                break
            elif count==1:
                break

            if media=='Print':
                rowValues = printRow(row, coverage)
            else:
                rowValues = onlineRow(row, coverage)

            table.cell(1,2).text_frame.text = rowValues[0] #nameplate
            table.cell(2,2).text_frame.text = rowValues[1] #title
            table.cell(3,2).text_frame.text = rowValues[2] #media
            date = rowValues[3].strftime("%Y-%m-%d")       #date
            table.cell(4,2).text_frame.text = date
            table.cell(5,2).text_frame.text = rowValues[4] #tone
            reach = str(rowValues[5])                      #reach
            table.cell(6,2).text_frame.text = reach
            ave = 'R ' + str(rowValues[6])                 #ave
            table.cell(7,2).text_frame.text = ave 
            table.cell(8,2).text_frame.text = rowValues[7] #link
            table.cell(9,2).text_frame.text = rowValues[8] #extract

            print('Values loaded')
            cell_list = populate_table_to_list(table)
            cell_list = set_table_styles(cell_list)
            table = set_table_styles_link_and_extract(table)
            
            count+=1

            if media == 'Print':
                LIST_PRINT_TOP10.append([rowValues[2], rowValues[1]])
            else:
                LIST_ONLINE_TOP10.append([rowValues[2], rowValues[1]])

    return table


def set_table_values(wb, ws1, prs, slide_list, coverage, media):

    #Colours of the highlighted cells
    colours = {'Product' : 'FFFFFF00',
               'Corporate' : 'FF92D050',
               'SSA Online': 'FFFFC000',
               'Social': 'FF00B0F0'}

    colour = colours[coverage]
    
    #creates list of rows from ws1
    rows = list(ws1.rows)

    #Count the number of clips added
    count = 0 

    #iterate first cell in each row until you find a highlighted cell 
    for row in rows:
        #If the cell is a certan colour
        if row[0].fill.fgColor.rgb == colour:
            print(row[0].value) 

            #Max 5 clips
            if (count==3) and (coverage=='Social'):
                break
            elif count==5:
                break

            if media=='Print':
                rowValues = printRow(row, coverage)
            else:
                rowValues = onlineRow(row, coverage)

            table = prs.slides[slide_list[count]]

            table.cell(1,2).text_frame.text = rowValues[0] #nameplate
            table.cell(2,2).text_frame.text = rowValues[1] #title
            table.cell(3,2).text_frame.text = rowValues[2] #media
            date = rowValues[3].strftime("%Y-%m-%d")       #date
            table.cell(4,2).text_frame.text = date
            table.cell(5,2).text_frame.text = rowValues[4] #tone
            reach = str(rowValues[5])                      #reach
            table.cell(6,2).text_frame.text = reach
            ave = 'R' + str(rowValues[6])                 #ave
            table.cell(7,2).text_frame.text = ave 
            table.cell(8,2).text_frame.text = rowValues[7] #link
            table.cell(9,2).text_frame.text = rowValues[8] #extract

            print('Values loaded')
            cell_list = populate_table_to_list(table)
            cell_list = set_table_styles(cell_list)
            table = set_table_styles_link_and_extract(table)
            
            count+=1

            if media == 'Print':
                LIST_PRINT_TOP10.append([rowValues[2], rowValues[1]])
            else:
                LIST_ONLINE_TOP10.append([rowValues[2], rowValues[1]])

    return prs

            


def top10Stories(wb, ws):
    #FIX!!!!!!!!!!!!!!!!!!
    """Generates Top10 from print and online clips
    """
    ws = ws['J11':'L20']

    #Randomly selects 5 print clips for print and online
    selection1 = random.sample(range(0, 10), 5)
    selection2 = random.sample(range(0, 10), 5)

    for index in range(5):
        if checkIfPrime(LIST_PRINT_TOP10[selection1[index]][0]):
            ws[index][2].value = 'PRIME'

        media = LIST_PRINT_TOP10[selection1[index]][0].replace(' [PRIME]', '')
        
        ws[index][0].value = media                                    #media
        ws[index][1].value = LIST_PRINT_TOP10[selection1[index]][1]   #title

    for index in range(5,10):
        if checkIfPrime(LIST_ONLINE_TOP10[selection2[index-5]][0]):
            ws[index][2].value = 'PRIME'

        media = LIST_ONLINE_TOP10[selection2[index-5]][0].replace(' [PRIME]',
                                                                  '')
        
        ws[index][0].value = media                                     #media
        ws[index][1].value = LIST_ONLINE_TOP10[selection2[index-5]][1] #title

    wb.save('done.xlsx')


start_time = time.time()

#PRINT_LIST, ONLINE_LIST, TV_LIST = loadPrime()

#  Reference to sheets
print_product_slides = [3,4,5,6,7]
print_coverage_slides = [15,16,17,18,19]

online_product_slides = [9,10,11,12,13]
online_coverage_slides = [21,22,23,24,25]

broadcast_slides = [27,28]

online_ssa_slides = [30,31,32,33,34]
social_ssa_slides = [36,37,38]

#  Load excel document
wb = load_workbook('AVE Data.xlsx', data_only=True)

#  Load Print, Online, Broadcas and SSA sheets
wsPrint = wb['Print']
wsOnline = wb['Online']
wsBroadcastList = wb['Broadcast']
wsSSA = wb['Online SSA']
wsSocial = wb['Social SSA']

# Load a presentation
prs = Presentation('data-files/Template.pptx')

# Get reference to second slide
sl = prs.slides[3]
table = sl.shapes[3].table

table = set_table_values_old(wb, wsPrint, table, 'Product', 'Print')
prs = set_table_values(wb, ws1, prs, print_product_slides, 'Product', 'Print')

prs.save('test.pptx')

print("Done!")
print("")
print("--- Took %s seconds ---" % (time.time() - start_time))
